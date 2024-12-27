from pypdf import PdfReader, PdfWriter
from pypdf.annotations import Text

from pptx import Presentation


class slideNotes():

    def __init__(self, pathToPPTx, pathToPDFtoChange,pathToPDFOutput=None):
        """
        Takes the original powerPoint presentation, the corresponding PDF
        and optionally the path for the output-File and puts the annotations
        written in the PowerPoint-Presentation as hidden notes to the PDF file.

        """
        self.pathToPPTx = pathToPPTx 
        self.pathToPDFtoChange = pathToPDFtoChange
        self.pathToPDFOutput = pathToPDFOutput

        ## Load the PPTX file ##
        # We store all annotations for each slide an an array
        self.prs = Presentation(pathToPPTx) # Load the presentation
        # Iterate through each slide
        self.slideNotesDict={}
        for index, slide in enumerate(self.prs.slides, start=1):
            # Get comments on the slide
            comment = slide.notes_slide.notes_text_frame.text    
            print(comment)
            self.slideNotesDict[index]=comment
        
        ## Init the PDF output ##
        self.reader = PdfReader(pathToPDFtoChange)
        self.writer = PdfWriter(clone_from=pathToPDFtoChange)
        # Get the document catalog 
        catalog = self.reader.trailer['/Root'] 
        # Get the page labels dictionary 
        if '/PageLabels' in catalog: 
            self.page_labels = catalog['/PageLabels']['/Nums'] 
        else: 
            self.page_labels = None
            print("Warning: There are no page-labels defined in the PDF. This means that either you do not "+
                  "have any animations, or you forgot to first run the script to add them to the PDF. Animations are declared with the PDF feature 'PageLabels'")
            assert self.reader.get_num_pages() == len(self.prs.slides),"PDF pages does align with the number of slides and no animations are defined within the PDF. Script will fail."


    
    def getNotesDict(self):
        return self.slideNotesDict
    
    def getAnnotation(self,slideNb):
        return self.slideNotesDict[slideNb]
    
    def getNumberOfPPTXSlides(self):
        return len(self.prs.slides)
    
    ## Write Output
    def addAnnotationToPDFPage(self,pdfPageNb, originalSlideNbFromPPTx):
        """
        Adds the annotations from the powerPointslide to the given pdfPageNb
        pageNb and slideNb starts at one
        """
        annotationText = self.getAnnotation(originalSlideNbFromPPTx)
        annotation = Text(
        text=annotationText,
        rect=(50, 550, 200, 650)
        )
        # Set annotation flags to 4 for printable annotations.
        # See "AnnotationFlag" for other options, e.g. hidden etc.
        #annotation.flags = 4
        # Seems like page_number starts at zero
        # eccob
        self.writer.add_annotation(page_number=pdfPageNb-1, annotation=annotation)
    
    def writeOutput(self):
        pathToStore = self.pathToPDFOutput
        if(self.pathToPDFOutput == None):
            print("Overwrite Original PDF")
            pathToStore=self.pathToPDFtoChange
        with open(pathToStore, "wb") as f:
            self.writer.write(f)
    
    def transferAnnotationsFromPPTxToPDF(self):
        """
        Writes the annotations from the PPP to the PDF.
        Won't overwrite already existing notes and as such, 
        duplicate entries may be the result if applied twice
        to the same document.
        """
        pdfPages = self.reader.get_num_pages()
        for pdfPage in range(pdfPages):
            pdfPageNb = pdfPage+1
            destPPTxPageNb = self.getPPTxSlideFromPDFPage(pdfPageNb)
            print(f"{pdfPageNb} - {destPPTxPageNb}")
            self.addAnnotationToPDFPage(pdfPageNb, destPPTxPageNb)

        
    # Function to get the page label for a specific page number
    def getPPTxSlideFromPDFPage(self,page_num):
        """
        Takes the PDF Page (1 to ...) and converts it back to the corresponding slide
        of the given PPTX. Expects that the PDF has the correct page-labels for that, 
        means the page label of each PDF corresponds to the correct slide of the PPTX.
        """
        correspondingPPTSlide = None
        if self.page_labels:
            for i in range(0, len(self.page_labels), 2):
                # We go through all page labels and find the fitting one for this page
                pdf_page = self.page_labels[i] + 1 # we start at one
                if (pdf_page <= page_num) and ( (i+2 >= len(self.page_labels)) or (page_num < (self.page_labels[i+2]+1)) ):
                    
                    label_dict = self.page_labels[i+1]
                    label_prefix = label_dict.get('/P', '')
                    label_start = label_dict.get('/St', 1)
                    label_pptSlide = label_start+(page_num - pdf_page)
                    print(f"{label_prefix}{label_pptSlide}")
                    correspondingPPTSlide = label_pptSlide
                    return correspondingPPTSlide
                #else:
                    # this is not a match, go on
            # This should not happen, but can
            raise ValueError("Page Labels are defined, but the PDF-labels are wrong")
        else:
            # No Page labels, no animations, no shift --> return the page_num
            
            return page_num

        
    


if __name__ == "__main__":
    inputPDFFileWithPageLabels = "example/ExampleAnimationsToPDF_withPageLabels.pdf"
    inputPPTXFile = "example/ExampleAnimationsToPDF_Original.pptx"
    outputPDFFile = "example/ExampleAnimationsToPDF_WithCommentsAndAnimations.pdf"
    slideNote = slideNotes(inputPPTXFile,inputPDFFileWithPageLabels,outputPDFFile)
    slideNote.transferAnnotationsFromPPTxToPDF()
    slideNote.writeOutput()
